from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BLUE = RGBColor(15, 40, 80)
BLUE = RGBColor(0, 120, 200)
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(30, 30, 30)
GRAY_TEXT = RGBColor(150, 150, 150)
LIGHT_BLUE_BG = RGBColor(240, 245, 252)
LIGHT_GREEN = RGBColor(232, 245, 233)
LIGHT_RED = RGBColor(255, 235, 238)

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, x, y, w, h, text, size, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tb

def header_bar(slide, title):
    add_rect(slide, 0, 0, prs.slide_width, Inches(1.2), DARK_BLUE)
    add_text(slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.8),
             title, Pt(28), bold=True, color=WHITE)
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(12.33), Pt(3), BLUE)
    add_rect(slide, 0, Inches(7.1), prs.slide_width, Inches(0.4), DARK_BLUE)

def footer(slide, text):
    add_text(slide, Inches(0.5), Inches(7.02), Inches(12.33), Inches(0.4),
             text, Pt(11), color=GRAY_TEXT, align=PP_ALIGN.RIGHT)

def bullets(slide, items, start_y=Inches(1.5), max_w=Inches(12)):
    tb = slide.shapes.add_textbox(Inches(0.7), start_y, max_w, Inches(5.3))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (icon, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = icon + "  " + text
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(10)

# ========== SLIDE 1: COVER ==========
s1 = blank_slide(prs)
add_rect(s1, 0, 0, prs.slide_width, prs.slide_height, DARK_BLUE)
add_text(s1, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5),
         "时光悦酩公寓底商市场调研报告", Pt(44), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, Inches(0.8), Inches(4.2), Inches(11.7), Inches(1),
         "浙江省杭州市上城区驿城路", Pt(22), color=RGBColor(180, 200, 230), align=PP_ALIGN.CENTER)
add_text(s1, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
         "2026年3月  |  出差汇报版", Pt(14), color=RGBColor(130, 160, 200), align=PP_ALIGN.CENTER)

# ========== SLIDE 2: TOC ==========
s2 = blank_slide(prs)
header_bar(s2, "报告目录")
items = [
    ("01", "项目基本情况"),
    ("02", "区域商圈分析"),
    ("03", "杭州市商业宏观数据"),
    ("04", "底商人流分析"),
    ("05", "业态组合建议"),
    ("06", "营商环境痛点"),
    ("07", "客户分层管理策略"),
    ("08", "综合研判与建议"),
]
bullets(s2, items, start_y=Inches(1.5))

# ========== SLIDE 3: BASIC INFO ==========
s3 = blank_slide(prs)
header_bar(s3, "项目基本情况")
items3 = [
    ("位置", "杭州市上城区驿城路，火车东站商圈核心辐射范围"),
    ("开发企业", "杭州弘驿置业有限公司"),
    ("物业公司", "杭州新天地高力国际物业管理有限公司"),
    ("总户数", "1961户  |  楼栋数：18栋"),
    ("绿化率", "30%  |  容积率：2.2"),
    ("物业费", "2.75-2.8元/m2/月"),
    ("二手房挂牌均价", "约38,727元/m2（2026年2月）"),
    ("定位", "住宅底商，服务本社区及周边短距离客群"),
]
bullets(s3, items3, start_y=Inches(1.5))
footer(s3, "数据来源：链家、贝壳")

# ========== SLIDE 4: LOCATION ==========
s4 = blank_slide(prs)
header_bar(s4, "区域商圈分析 - 火车东站商圈")
items4 = [
    ("交通枢纽", "杭州东站系长三角核心枢纽，日均客流庞大，带来过境消费人群"),
    ("住宅密度", "周边多个成熟小区（天聚府、明月嘉苑等），常住人口基数大"),
    ("商业配套", "1公里内有弄口滨河汇生活广场、江南晟世生活广场"),
    ("竞争压力", "区域内住宅底商体量大，同质化竞争激烈"),
    ("地铁优势", "紧邻地铁9号线驿城路站，步行约350米"),
]
bullets(s4, items4, start_y=Inches(1.5))

# ========== SLIDE 5: MARKET DATA (KPI cards) ==========
s5 = blank_slide(prs)
header_bar(s5, "杭州市零售商业市场数据（2025年）")
kpis = [
    ("全市优质商业存量", "1775.3", "万m2"),
    ("2025年新增供应", "105", "万m2"),
    ("优质商业平均租金", "6.48", "元/天/m2  |  较去年↓3.4%"),
    ("全市平均空置率", "9.7%", "较去年↑0.8pp"),
    ("城东新城空置率", ">15%", "非核心商圈  |  需重点关注"),
]
n = len(kpis)
card_w = Inches(11.5) / n
start_x = Inches(0.9)
for i, (label, value, unit) in enumerate(kpis):
    x = start_x + i * card_w
    add_rect(s5, x, Inches(1.8), card_w - Inches(0.15), Inches(3.0), LIGHT_BLUE_BG, RGBColor(200,220,240))
    add_text(s5, x + Inches(0.1), Inches(2.0), card_w - Inches(0.3), Inches(1.0),
             value, Pt(36), bold=True, color=RGBColor(0,100,180), align=PP_ALIGN.CENTER)
    add_text(s5, x + Inches(0.1), Inches(3.1), card_w - Inches(0.3), Inches(0.4),
             unit, Pt(12), color=RGBColor(100,130,160), align=PP_ALIGN.CENTER)
    add_text(s5, x + Inches(0.1), Inches(3.55), card_w - Inches(0.3), Inches(1.1),
             label, Pt(13), color=RGBColor(60,60,60), align=PP_ALIGN.CENTER)
footer(s5, "数据来源：戴德梁行《2025年杭州商业地产报告》")

# ========== SLIDE 6: FOOT TRAFFIC ==========
s6 = blank_slide(prs)
header_bar(s6, "底商人流分析")
items6 = [
    ("社区居民", "1961户，约5000-6000人 | 中等消费，日常刚需为主"),
    ("周边居民", "步行5-10分钟可达 | 中等偏低消费"),
    ("地铁客流", "驿城路站通勤人群 | 随机性消费，便利性需求"),
    ("东站转乘", "出差/过境人群 | 低频，价格敏感"),
    ("峰值时段", "早7-9点（早餐）| 午11:30-13:30（快餐）| 晚17:30-20:00（生活服务）"),
    ("人流短板", "距火车东站主入口约800米，过境客流转化率有限"),
]
bullets(s6, items6, start_y=Inches(1.5))
footer(s6, "注：人流数据为公开信息推断，建议实地调研核实")

# ========== SLIDE 7: BUSINESS MIX ==========
s7 = blank_slide(prs)
header_bar(s7, "业态组合建议")
cols = [
    ("必须配套", ["便利店/生鲜超市", "早餐/快餐"]),
    ("重要业态", ["快递驿站", "美容美发/美甲", "社区教育培训"]),
    ("可选/慎选", ["宠物店（可选）", "棋牌室（慎选）", "重油烟餐饮（禁止）"]),
]
col_w = Inches(3.9)
gap = Inches(0.3)
for i, (hdr, items) in enumerate(cols):
    x = Inches(0.5) + i * (col_w + gap)
    add_rect(s7, x, Inches(1.5), col_w, Inches(0.55), BLUE)
    add_text(s7, x, Inches(1.55), col_w, Inches(0.5),
             hdr, Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb = s7.shapes.add_textbox(x, Inches(2.2), col_w, Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "* " + item
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)

# ========== SLIDE 8: PAIN POINTS ==========
s8 = blank_slide(prs)
header_bar(s8, "营商环境痛点")
items8 = [
    ("同质化", "周边底商以餐饮为主，差异化不足"),
    ("消费力", "区域以刚需客群为主，客单价提升空间有限"),
    ("电商冲击", "便利店+餐饮受线上配送分流，线下体验型业态更具优势"),
    ("停车难", "住宅底商停车位普遍不足，制约大宗消费"),
]
tb8 = s8.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.0))
tf8 = tb8.text_frame
tf8.word_wrap = True
for i, (title, desc) in enumerate(items8):
    p = tf8.paragraphs[0] if i == 0 else tf8.add_paragraph()
    p.text = "[!] " + title + "：" + desc
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(180, 50, 50)
    p.space_after = Pt(16)

# ========== SLIDE 9: TENANT MANAGEMENT ==========
s9 = blank_slide(prs)
header_bar(s9, "客户分层管理模型（ABC分类）")
cols9 = [
    ("A类 - 核心客户", ["品牌连锁/现金流稳定", "租约3年以上", "免租装修期优惠", "重点维护关系"]),
    ("B类 - 成长客户", ["运营稳定/面积中等", "每年谈递增", "保持沟通频率", "观察培养潜力"]),
    ("C类 - 淘汰客户", ["频繁欠租", "业态不符合定位", "到期不续约", "提前储备替代品牌"]),
]
for i, (hdr, items) in enumerate(cols9):
    x = Inches(0.5) + i * Inches(4.2)
    colors = [RGBColor(0,100,180), RGBColor(200,150,0), RGBColor(180,60,60)]
    add_rect(s9, x, Inches(1.5), Inches(3.9), Inches(0.55), colors[i])
    add_text(s9, x, Inches(1.55), Inches(3.9), Inches(0.5),
             hdr, Pt(17), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb = s9.shapes.add_textbox(x, Inches(2.2), Inches(3.9), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "- " + item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)

# ========== SLIDE 10: TENANT CARE ==========
s10 = blank_slide(prs)
header_bar(s10, "客户维护具体措施")
items10 = [
    ("日常维护", "建立租户微信群，节假日问候+活动推送；每季度上门拜访，了解经营状况"),
    ("物业服务", "及时响应报修需求，服务响应速度是核心竞争力"),
    ("租约管理", "主力租户（>100m2）签约不超过3年，年度递增3%-5%；小型租户（<50m2）签1+1模式"),
    ("风险预警", "逾期15天以上主动介入，提前化解欠租风险"),
    ("业态管控", "禁止重油烟餐饮、棋牌室；鼓励轻餐饮、烘焙、宠物、教育培训"),
]
bullets(s10, items10, start_y=Inches(1.5))

# ========== SLIDE 11: SUMMARY ==========
s11 = blank_slide(prs)
header_bar(s11, "综合研判与建议")
# Left: Opportunity
add_rect(s11, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.0), LIGHT_GREEN, RGBColor(76,175,80))
add_text(s11, Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.5),
         "[+] 机会", Pt(20), bold=True, color=RGBColor(27,94,32))
opp_items = [
    "地铁9号线+1961户社区居民，稳定基础客群",
    "社区配套型业态（便利店/快递）抗电商冲击能力强",
    "空置率偏高，竞争洗牌后优质底商更稀缺",
]
tb_opp = s11.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.0))
tf_opp = tb_opp.text_frame
tf_opp.word_wrap = True
for i, t in enumerate(opp_items):
    p = tf_opp.paragraphs[0] if i == 0 else tf_opp.add_paragraph()
    p.text = "- " + t
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(30,60,30)
    p.space_after = Pt(10)
# Right: Risk
add_rect(s11, Inches(6.9), Inches(1.5), Inches(5.9), Inches(5.0), LIGHT_RED, RGBColor(229,57,53))
add_text(s11, Inches(7.1), Inches(1.6), Inches(5.5), Inches(0.5),
         "[-] 风险", Pt(20), bold=True, color=RGBColor(183,28,28))
risk_items = [
    "城东新城非核心商圈，租金上行空间有限",
    "区域内底商同质化严重，差异化能力是关键",
    "消费力分流，需精准定位客群",
]
tb_risk = s11.shapes.add_textbox(Inches(7.1), Inches(2.2), Inches(5.5), Inches(4.0))
tf_risk = tb_risk.text_frame
tf_risk.word_wrap = True
for i, t in enumerate(risk_items):
    p = tf_risk.paragraphs[0] if i == 0 else tf_risk.add_paragraph()
    p.text = "- " + t
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(60,30,30)
    p.space_after = Pt(10)
# Bottom advice bar
add_rect(s11, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.5), DARK_BLUE)
advice = "核心建议：业态优先便民+高频 | ABC租户分级管理 | 利用地铁人流峰值导流 | 控制单店面积抗风险"
add_text(s11, Inches(0.5), Inches(6.62), Inches(12.33), Inches(0.45),
         advice, Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ========== SLIDE 12: BACK COVER ==========
s12 = blank_slide(prs)
add_rect(s12, 0, 0, prs.slide_width, prs.slide_height, DARK_BLUE)
add_text(s12, Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.5),
         "谢谢", Pt(52), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s12, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.8),
         "时光悦酩公寓底商市场调研报告", Pt(20), color=RGBColor(180,200,230), align=PP_ALIGN.CENTER)
add_text(s12, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
         "2026年3月", Pt(14), color=RGBColor(130,160,200), align=PP_ALIGN.CENTER)

# SAVE
out = r"E:\openclaw-work\slides\hangzhou-commercial-report\hangzhou-commercial-report-clean.pptx"
prs.save(out)
import os
size_kb = os.path.getsize(out) // 1024
print("Done:", out)
print("Size:", size_kb, "KB")
